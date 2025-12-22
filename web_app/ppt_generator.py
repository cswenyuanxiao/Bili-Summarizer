from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import io
import re

class PPTGenerator:
    def __init__(self):
        self.prs = Presentation()
        # 16:9 Aspect Ratio
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    def _set_font(self, run, size: int, bold: bool = False, color: tuple = None):
        font = run.font
        font.name = 'Microsoft YaHei'  # Fallback to standard font
        font.size = Pt(size)
        font.bold = bold
        if color:
            font.color.rgb = RGBColor(*color)

    def create_cover_slide(self, title: str, subtitle: str):
        """Create the title slide."""
        slide_layout = self.prs.slide_layouts[0] # Title Slide
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title
        title_shape = slide.shapes.title
        title_shape.text = title
        self._set_font(title_shape.text_frame.paragraphs[0].runs[0], 44, True, (31, 41, 55)) # Dark Grey

        # Subtitle
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle
        self._set_font(subtitle_shape.text_frame.paragraphs[0].runs[0], 24, False, (107, 114, 128)) # Light Grey

    def create_content_slide(self, title: str, content_points: list):
        """Create a standard content slide with bullet points."""
        slide_layout = self.prs.slide_layouts[1] # Title and Content
        slide = self.prs.slides.add_slide(slide_layout)

        # Title
        title_shape = slide.shapes.title
        title_shape.text = title
        self._set_font(title_shape.text_frame.paragraphs[0].runs[0], 36, True, (0, 0, 0))

        # Content
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear() # Clear default empty paragraph

        for point in content_points:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            # Check if point looks like a sub-point (starts with - or space)
            if point.startswith('  ') or point.startswith('\t'):
                p.level = 1
                p.text = point.strip()
            
            self._set_font(p.runs[0], 20, False, (55, 65, 81))
            p.space_after = Pt(10)

    def generate_from_json(self, data: dict) -> io.BytesIO:
        """
        Generate PPT from structured JSON data.
        Data format:
        {
            "title": "Main Title",
            "subtitle": "Author/Source",
            "slides": [
                {
                    "title": "Slide Title",
                    "content": ["Point 1", "Point 2"]
                },
                ...
            ]
        }
        """
        # 1. Cover
        self.create_cover_slide(data.get('title', '视频总结'), data.get('subtitle', 'AI Auto-Generated'))

        # 2. Content Slides
        for slide_data in data.get('slides', []):
            self.create_content_slide(slide_data.get('title', 'Unknown'), slide_data.get('content', []))

        # 3. Save to buffer
        output = io.BytesIO()
        self.prs.save(output)
        output.seek(0)
        return output
